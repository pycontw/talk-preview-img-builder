"""
ImgBuilder Class
"""
import os
from dataclasses import dataclass
from pathlib import Path
from typing import List, Tuple

from constant import CATEGORY_TO_TEXT, LANGUAGE_TO_TEXT, PYTHON_LEVEL_TO_TEXT
from loguru import logger
from model import Material, Speaker
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Cm, Pt


@dataclass
class TalkPreviewPPTBuilder:  # pylint: disable=too-many-instance-attributes
    """
    PPTBuilder for building image by given materials
    """

    material: Material
    size: Tuple[float, float]
    title_upper_left_pos: Tuple[float, float]
    content_upper_left_pos: Tuple[float, float]
    footer_upper_left_pos: Tuple[float, float]
    title_height: float
    content_height: float
    # TODO: Update  to use left position and also need to unify interface of img builder
    speaker_right_pos: float
    speaker_upper_pos: float
    speaker_width: float = 10
    speaker_margin: float = 1.6
    footer_width: float = 12
    font: str = "PingFang.ttc"
    bold_font: str = "PingFang.ttc"
    text_color: str = "#000000"
    hightlight_text_color: str = "#ffffff"
    output_path: Path = Path(__file__).parent.parent / "export" / "ppt"
    file_name: str = "talk_preview.pptx"

    def __post_init__(self):
        """
        Post initialization
        """
        self.content_size = (
            self.size[0] - self.content_upper_left_pos[0] * 2,
            self.content_height,
        )
        self.title_size: Tuple[float, float] = (
            self.size[0] - self.title_upper_left_pos[0] * 2,
            self.title_height,
        )

    def __fill_in_title(self, slide: Slide, title: str):
        """
        Fill in the title
        """
        text_box = slide.shapes.add_textbox(
            Cm(self.title_upper_left_pos[0]),
            Cm(self.title_upper_left_pos[1]),
            Cm(self.title_size[0]),
            Cm(self.title_size[1]),
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        paragraph = text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.CENTER
        logger.debug("Filling title {}", title)
        paragraph.text = title

        paragraph.font.size = Pt(20)
        paragraph.font.name = self.font
        hex_code = self.text_color.lstrip("#")
        paragraph.font.color.rgb = RGBColor(
            *tuple(int(hex_code[i : i + 2], 16) for i in (0, 2, 4))
        )
        paragraph.font.bold = True

    def __fill_in_abstract(self, slide: Slide, abstract: str):
        """
        Fill in the abstract of tha talk
        """
        text_box = slide.shapes.add_textbox(
            Cm(self.content_upper_left_pos[0]),
            Cm(self.content_upper_left_pos[1]),
            Cm(self.content_size[0]),
            Cm(self.content_size[1]),
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        paragraph = text_frame.add_paragraph()
        abstract = abstract.replace("\r", "")
        logger.debug("Filling abstract {}", abstract)
        paragraph.text = abstract

        paragraph.line_spacing = Pt(12)
        paragraph.font.size = Pt(10)
        paragraph.font.name = self.font
        hex_code = self.text_color.lstrip("#")
        paragraph.font.color.rgb = RGBColor(
            *tuple(int(hex_code[i : i + 2], 16) for i in (0, 2, 4))
        )

    def __fill_in_speaker(self, slide: Slide, speakers: List[Speaker]):
        """
        Fill in the speaker of the talk
        """
        text_box = slide.shapes.add_textbox(
            Cm(
                self.speaker_right_pos - self.speaker_width - self.speaker_margin
            ),  # TODO: Update to use left position
            Cm(self.speaker_upper_pos),
            Cm(self.speaker_width),
            Cm(1),
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        paragraph = text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.RIGHT
        name_list = "、".join([speaker.name for speaker in speakers])
        display_text = f"— {name_list}"
        logger.debug("Filling speakers {}", display_text)
        paragraph.text = display_text

        paragraph.font.size = Pt(16)
        paragraph.font.name = self.font
        hex_code = self.text_color.lstrip("#")
        paragraph.font.color.rgb = RGBColor(
            *tuple(int(hex_code[i : i + 2], 16) for i in (0, 2, 4))
        )

    def __fill_in_footer(
        self, slide: Slide, category: str, python_level: str, language: str
    ):
        """
        Fill in the category, python level, language of the talk
        """
        display_text = "    .    ".join(
            [
                CATEGORY_TO_TEXT[category],
                f"語言：{LANGUAGE_TO_TEXT[language]}",
                f"Python 難易度：{PYTHON_LEVEL_TO_TEXT[python_level]}",
            ]
        )

        text_box = slide.shapes.add_textbox(
            Cm(self.footer_upper_left_pos[0]),
            Cm(self.footer_upper_left_pos[1]),
            Cm(self.footer_width),
            Cm(1),
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        paragraph = text_frame.add_paragraph()
        logger.debug("Filling footer {}", display_text)
        paragraph.text = display_text

        paragraph.font.size = Pt(9)
        paragraph.font.name = self.font
        hex_code = self.hightlight_text_color.lstrip("#")
        paragraph.font.color.rgb = RGBColor(
            *tuple(int(hex_code[i : i + 2], 16) for i in (0, 2, 4))
        )

    def execute(self):
        """
        Execute the builder
        """
        if not self.material.background_img:
            return
        if not self.output_path.exists():
            os.mkdir(self.output_path.resolve())
        presentation = Presentation()
        for talk in self.material.talk_list:
            logger.info("Building slide for speech: {}...", talk.title)
            blank_slide_layout = presentation.slide_layouts[6]
            slide = presentation.slides.add_slide(blank_slide_layout)
            left = top = Cm(0)
            slide.shapes.add_picture(
                str(self.material.background_img_path.resolve()),
                left,
                top,
                width=Cm(self.size[0]),
                height=Cm(self.size[1]),
            )
            self.__fill_in_title(slide, talk.title)
            self.__fill_in_abstract(slide, talk.abstract)
            self.__fill_in_speaker(slide, talk.speakers)
            self.__fill_in_footer(
                slide,
                category=talk.category,
                python_level=talk.python_level,
                language=talk.language,
            )
        presentation.save((self.output_path / self.file_name).resolve())
